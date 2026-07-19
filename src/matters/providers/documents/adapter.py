"""Deterministic document extraction behind safe, injectable extractors."""

from __future__ import annotations

from collections.abc import Mapping, Sequence
from dataclasses import dataclass, field
from io import BytesIO
from typing import Protocol


DOCUMENT_STATUSES = frozenset(
    {
        "extracted",
        "partial",
        "not_tracked",
        "unsupported",
        "encrypted",
        "corrupt",
        "resource_exhausted",
    }
)
ANCHOR_REQUIREMENTS = {
    "line": frozenset({"line_start", "line_end", "char_start", "char_end"}),
    "page": frozenset({"page"}),
    "passage": frozenset({"page", "char_start", "char_end"}),
    "slide": frozenset({"slide"}),
    "shape": frozenset({"slide", "shape"}),
    "sheet_cell": frozenset({"sheet", "cell"}),
    "sheet_range": frozenset({"sheet", "range"}),
}
TEXT_MEDIA_TYPES = frozenset(
    {
        "text/plain",
        "text/markdown",
        "text/x-python",
        "text/x-source",
        "application/json",
        "application/xml",
        "application/yaml",
        "application/toml",
    }
)


@dataclass(frozen=True)
class ExtractionBudget:
    max_bytes: int = 16 * 1024 * 1024
    max_parts: int = 100_000
    max_characters: int = 16 * 1024 * 1024

    def __post_init__(self) -> None:
        if min(self.max_bytes, self.max_parts, self.max_characters) < 1:
            raise ValueError("document extraction budgets must be positive")


@dataclass(frozen=True)
class DocumentSource:
    source_version_id: str
    external_id: str
    media_type: str
    content: bytes
    tracking_disposition: str
    parent_source_version_id: str = ""
    metadata: Mapping[str, object] = field(default_factory=dict)

    def __post_init__(self) -> None:
        if not self.source_version_id or not self.external_id or not self.media_type:
            raise ValueError(
                "source_version_id, external_id, and media_type are required"
            )
        object.__setattr__(self, "content", bytes(self.content))
        object.__setattr__(self, "metadata", dict(self.metadata))


@dataclass(frozen=True)
class DocumentAnchor:
    kind: str
    location: Mapping[str, object]
    text: str
    modality: str = "observed"
    confidence: float = 1.0

    def __post_init__(self) -> None:
        required = ANCHOR_REQUIREMENTS.get(self.kind)
        if required is None:
            raise ValueError(f"unsupported document anchor kind: {self.kind}")
        location = dict(self.location)
        missing = required - set(location)
        if missing:
            raise ValueError(
                f"document anchor {self.kind} missing: {sorted(missing)}"
            )
        if not 0.0 <= self.confidence <= 1.0:
            raise ValueError("document anchor confidence must be between 0 and 1")
        object.__setattr__(self, "location", location)


@dataclass(frozen=True)
class DocumentExtractionResult:
    status: str
    source_version_id: str
    extractor_id: str
    anchors: tuple[DocumentAnchor, ...] = ()
    gaps: tuple[str, ...] = ()
    parent_source_version_id: str = ""

    def __post_init__(self) -> None:
        if self.status not in DOCUMENT_STATUSES:
            raise ValueError(f"unsupported document status: {self.status}")
        object.__setattr__(self, "anchors", tuple(self.anchors))
        object.__setattr__(self, "gaps", tuple(self.gaps))


class DocumentExtractor(Protocol):
    """A safe extractor supplied by the runtime, never discovered content."""

    extractor_id: str
    media_types: frozenset[str]
    executes_active_content: bool
    activates_external_data: bool

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        """Return bounded, parent-bound parts and visible gaps."""


class PlainTextExtractor:
    """UTF-8-only line/character extraction with no format execution."""

    extractor_id = "matters.plain-text:utf8-lines:v1"
    media_types = TEXT_MEDIA_TYPES
    executes_active_content = False
    activates_external_data = False

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        if len(source.content) > budget.max_bytes:
            return DocumentExtractionResult(
                "resource_exhausted",
                source.source_version_id,
                self.extractor_id,
                gaps=("document_byte_budget_exceeded",),
                parent_source_version_id=source.parent_source_version_id,
            )
        try:
            text = source.content.decode("utf-8-sig", errors="strict")
        except UnicodeDecodeError:
            return DocumentExtractionResult(
                "unsupported",
                source.source_version_id,
                self.extractor_id,
                gaps=("utf8_decode_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        if len(text) > budget.max_characters:
            return DocumentExtractionResult(
                "resource_exhausted",
                source.source_version_id,
                self.extractor_id,
                gaps=("document_character_budget_exceeded",),
                parent_source_version_id=source.parent_source_version_id,
            )

        anchors: list[DocumentAnchor] = []
        character_offset = 0
        for line_number, line_with_ending in enumerate(
            text.splitlines(keepends=True),
            start=1,
        ):
            line = line_with_ending.rstrip("\r\n")
            line_end = character_offset + len(line)
            anchors.append(
                DocumentAnchor(
                    "line",
                    {
                        "line_start": line_number,
                        "line_end": line_number,
                        "char_start": character_offset,
                        "char_end": line_end,
                    },
                    line,
                )
            )
            character_offset += len(line_with_ending)
            if len(anchors) == budget.max_parts:
                remaining = line_number < len(text.splitlines())
                return DocumentExtractionResult(
                    "partial" if remaining else "extracted",
                    source.source_version_id,
                    self.extractor_id,
                    anchors=tuple(anchors),
                    gaps=(
                        ("document_part_budget_exceeded",) if remaining else ()
                    ),
                    parent_source_version_id=source.parent_source_version_id,
                )
        if text and not anchors:
            anchors.append(
                DocumentAnchor(
                    "line",
                    {
                        "line_start": 1,
                        "line_end": 1,
                        "char_start": 0,
                        "char_end": len(text),
                    },
                    text,
                )
            )
        return DocumentExtractionResult(
            "extracted",
            source.source_version_id,
            self.extractor_id,
            anchors=tuple(anchors),
            parent_source_version_id=source.parent_source_version_id,
        )


class PdfTextExtractor:
    extractor_id = "matters.pdf:pymupdf-text:v1"
    media_types = frozenset({"application/pdf"})
    executes_active_content = False
    activates_external_data = False

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        if len(source.content) > budget.max_bytes:
            return DocumentExtractionResult(
                "resource_exhausted",
                source.source_version_id,
                self.extractor_id,
                gaps=("document_byte_budget_exceeded",),
                parent_source_version_id=source.parent_source_version_id,
            )
        try:
            import fitz

            document = fitz.open(stream=source.content, filetype="pdf")
        except Exception:
            return DocumentExtractionResult(
                "corrupt",
                source.source_version_id,
                self.extractor_id,
                gaps=("pdf_open_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        if bool(document.needs_pass):
            document.close()
            return DocumentExtractionResult(
                "encrypted",
                source.source_version_id,
                self.extractor_id,
                gaps=("pdf_password_required",),
                parent_source_version_id=source.parent_source_version_id,
            )
        anchors: list[DocumentAnchor] = []
        characters = 0
        partial = False
        try:
            for page_index, page in enumerate(document, start=1):
                text = page.get_text("text").strip()
                if not text:
                    continue
                characters += len(text)
                if (
                    characters > budget.max_characters
                    or len(anchors) >= budget.max_parts
                ):
                    partial = True
                    break
                anchors.append(
                    DocumentAnchor(
                        "page",
                        {"page": page_index},
                        text,
                    )
                )
        finally:
            document.close()
        return DocumentExtractionResult(
            "partial" if partial else "extracted",
            source.source_version_id,
            self.extractor_id,
            anchors=tuple(anchors),
            gaps=(("document_part_budget_exceeded",) if partial else ()),
            parent_source_version_id=source.parent_source_version_id,
        )


class DocxTextExtractor:
    extractor_id = "matters.docx:python-docx-text:v1"
    media_types = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        }
    )
    executes_active_content = False
    activates_external_data = False

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        if len(source.content) > budget.max_bytes:
            return DocumentExtractionResult(
                "resource_exhausted",
                source.source_version_id,
                self.extractor_id,
                gaps=("document_byte_budget_exceeded",),
                parent_source_version_id=source.parent_source_version_id,
            )
        try:
            from docx import Document

            document = Document(BytesIO(source.content))
        except Exception:
            return DocumentExtractionResult(
                "corrupt",
                source.source_version_id,
                self.extractor_id,
                gaps=("docx_open_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        anchors: list[DocumentAnchor] = []
        offset = 0
        characters = 0
        for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if not text:
                continue
            characters += len(text)
            if characters > budget.max_characters or len(anchors) >= budget.max_parts:
                return DocumentExtractionResult(
                    "partial",
                    source.source_version_id,
                    self.extractor_id,
                    anchors=tuple(anchors),
                    gaps=("document_part_budget_exceeded",),
                    parent_source_version_id=source.parent_source_version_id,
                )
            anchors.append(
                DocumentAnchor(
                    "line",
                    {
                        "line_start": paragraph_index,
                        "line_end": paragraph_index,
                        "char_start": offset,
                        "char_end": offset + len(text),
                    },
                    text,
                )
            )
            offset += len(text) + 1
        return DocumentExtractionResult(
            "extracted",
            source.source_version_id,
            self.extractor_id,
            anchors=tuple(anchors),
            parent_source_version_id=source.parent_source_version_id,
        )


class PptxTextExtractor:
    extractor_id = "matters.pptx:python-pptx-text:v1"
    media_types = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        }
    )
    executes_active_content = False
    activates_external_data = False

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        try:
            from pptx import Presentation

            presentation = Presentation(BytesIO(source.content))
        except Exception:
            return DocumentExtractionResult(
                "corrupt",
                source.source_version_id,
                self.extractor_id,
                gaps=("pptx_open_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        anchors: list[DocumentAnchor] = []
        characters = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = str(getattr(shape, "text", "")).strip()
                if not text:
                    continue
                characters += len(text)
                if characters > budget.max_characters or len(anchors) >= budget.max_parts:
                    return DocumentExtractionResult(
                        "partial",
                        source.source_version_id,
                        self.extractor_id,
                        anchors=tuple(anchors),
                        gaps=("document_part_budget_exceeded",),
                        parent_source_version_id=source.parent_source_version_id,
                    )
                anchors.append(
                    DocumentAnchor(
                        "shape",
                        {"slide": slide_index, "shape": shape_index},
                        text,
                    )
                )
        return DocumentExtractionResult(
            "extracted",
            source.source_version_id,
            self.extractor_id,
            anchors=tuple(anchors),
            parent_source_version_id=source.parent_source_version_id,
        )


class XlsxValueExtractor:
    extractor_id = "matters.xlsx:openpyxl-values:v1"
    media_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}
    )
    executes_active_content = False
    activates_external_data = False

    def extract(
        self,
        source: DocumentSource,
        *,
        budget: ExtractionBudget,
    ) -> DocumentExtractionResult:
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(
                BytesIO(source.content),
                read_only=True,
                data_only=True,
                keep_links=False,
            )
        except Exception:
            return DocumentExtractionResult(
                "corrupt",
                source.source_version_id,
                self.extractor_id,
                gaps=("xlsx_open_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        anchors: list[DocumentAnchor] = []
        characters = 0
        try:
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        text = str(cell.value)
                        characters += len(text)
                        if (
                            characters > budget.max_characters
                            or len(anchors) >= budget.max_parts
                        ):
                            return DocumentExtractionResult(
                                "partial",
                                source.source_version_id,
                                self.extractor_id,
                                anchors=tuple(anchors),
                                gaps=("document_part_budget_exceeded",),
                                parent_source_version_id=source.parent_source_version_id,
                            )
                        anchors.append(
                            DocumentAnchor(
                                "sheet_cell",
                                {"sheet": sheet.title, "cell": cell.coordinate},
                                text,
                            )
                        )
        finally:
            workbook.close()
        return DocumentExtractionResult(
            "extracted",
            source.source_version_id,
            self.extractor_id,
            anchors=tuple(anchors),
            parent_source_version_id=source.parent_source_version_id,
        )


class DocumentAdapter:
    """Select one declared extractor and enforce the non-execution boundary."""

    provider_id = "documents"

    def __init__(
        self,
        extractors: Sequence[DocumentExtractor] = (),
        *,
        budget: ExtractionBudget | None = None,
    ):
        self._budget = budget or ExtractionBudget()
        candidates: dict[str, DocumentExtractor] = {}
        for extractor in (
            PlainTextExtractor(),
            PdfTextExtractor(),
            DocxTextExtractor(),
            PptxTextExtractor(),
            XlsxValueExtractor(),
            *tuple(extractors),
        ):
            if (
                extractor.executes_active_content
                or extractor.activates_external_data
            ):
                raise ValueError(
                    "document extractors may not execute active content "
                    "or activate external data"
                )
            if not extractor.extractor_id:
                raise ValueError("document extractor identity is required")
            for media_type in extractor.media_types:
                if media_type in candidates:
                    raise ValueError(
                        f"multiple document extractors own {media_type}"
                    )
                candidates[media_type] = extractor
        self._extractors = candidates

    def capability(self, media_type: str) -> Mapping[str, object]:
        extractor = self._extractors.get(media_type)
        if extractor is None:
            return {
                "status": "unsupported",
                "media_type": media_type,
                "extractor_id": "",
                "active_content_execution": False,
                "external_data_activation": False,
            }
        return {
            "status": "configured",
            "media_type": media_type,
            "extractor_id": extractor.extractor_id,
            "active_content_execution": False,
            "external_data_activation": False,
        }

    def extract(self, source: DocumentSource) -> DocumentExtractionResult:
        if source.tracking_disposition != "tracked":
            return DocumentExtractionResult(
                "not_tracked",
                source.source_version_id,
                "",
                gaps=("current_tracked_disposition_required",),
                parent_source_version_id=source.parent_source_version_id,
            )
        extractor = self._extractors.get(source.media_type)
        if extractor is None:
            return DocumentExtractionResult(
                "unsupported",
                source.source_version_id,
                "",
                gaps=("safe_extractor_unavailable",),
                parent_source_version_id=source.parent_source_version_id,
            )
        try:
            result = extractor.extract(source, budget=self._budget)
        except Exception:
            return DocumentExtractionResult(
                "corrupt",
                source.source_version_id,
                extractor.extractor_id,
                gaps=("document_extractor_failed",),
                parent_source_version_id=source.parent_source_version_id,
            )
        if result.source_version_id != source.source_version_id:
            raise ValueError("document extractor returned a foreign source version")
        if result.extractor_id != extractor.extractor_id:
            raise ValueError("document extractor identity changed during extraction")
        if (
            result.parent_source_version_id
            != source.parent_source_version_id
        ):
            raise ValueError(
                "document extractor returned a foreign parent source version"
            )
        if len(result.anchors) > self._budget.max_parts:
            raise ValueError("document extractor exceeded the declared part budget")
        for anchor in result.anchors:
            if len(anchor.text) > self._budget.max_characters:
                raise ValueError(
                    "document extractor returned an oversized anchor"
                )
        return result


__all__ = [
    "ANCHOR_REQUIREMENTS",
    "DOCUMENT_STATUSES",
    "DocumentAdapter",
    "DocumentAnchor",
    "DocumentExtractionResult",
    "DocumentExtractor",
    "DocumentSource",
    "DocxTextExtractor",
    "ExtractionBudget",
    "PlainTextExtractor",
    "PdfTextExtractor",
    "PptxTextExtractor",
    "TEXT_MEDIA_TYPES",
    "XlsxValueExtractor",
]
